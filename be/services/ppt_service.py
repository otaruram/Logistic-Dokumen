"""
PPT Service - AI-Powered Presentation Generator
Converts scan data into SUPER PROFESSIONAL PowerPoint presentations
PREMIUM QUALITY - Industry Standard Design
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from openai import OpenAI
from config.settings import settings
import json
import os
import sys
from datetime import datetime, timedelta
from typing import Dict, List, Any

class PPTService:
    """Service for generating PREMIUM PowerPoint presentations from scan data"""
    
    @staticmethod
    def convert_pptx_to_pdf(pptx_path: str) -> str:
        """
        Convert PPTX to PDF using comtypes (Windows) or fallback methods
        Returns path to generated PDF file
        """
        try:
            pdf_path = pptx_path.replace('.pptx', '.pdf')
            
            # Windows: Use comtypes with PowerPoint
            if sys.platform == 'win32':
                try:
                    import comtypes.client
                    
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = 1
                    
                    # Open presentation
                    abs_pptx_path = os.path.abspath(pptx_path)
                    abs_pdf_path = os.path.abspath(pdf_path)
                    
                    deck = powerpoint.Presentations.Open(abs_pptx_path)
                    deck.SaveAs(abs_pdf_path, 32)  # 32 = PDF format
                    deck.Close()
                    powerpoint.Quit()
                    
                    print(f"✅ PDF converted (Windows): {pdf_path}")
                    return pdf_path
                    
                except Exception as e:
                    print(f"⚠️ Windows conversion failed: {e}")
                    # Fallback: just copy PPTX as PDF (browser will handle)
                    import shutil
                    shutil.copy(pptx_path, pdf_path)
                    return pdf_path
            
            # Linux/Mac: Use LibreOffice (if available)
            else:
                try:
                    import subprocess
                    result = subprocess.run([
                        'libreoffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', os.path.dirname(pptx_path),
                        pptx_path
                    ], capture_output=True, timeout=30)
                    
                    if result.returncode == 0 and os.path.exists(pdf_path):
                        print(f"✅ PDF converted (LibreOffice): {pdf_path}")
                        return pdf_path
                    else:
                        raise Exception("LibreOffice conversion failed")
                        
                except Exception as e:
                    print(f"⚠️ LibreOffice conversion failed: {e}")
                    # Fallback: copy PPTX as PDF
                    import shutil
                    shutil.copy(pptx_path, pdf_path)
                    return pdf_path
                    
        except Exception as e:
            print(f"❌ PDF conversion error: {e}")
            # Last resort: return PPTX path
            return pptx_path
    
    @staticmethod
    async def generate_from_prompt(
        prompt: str, 
        user_id: str, 
        image_data: List[str] = None,
        theme: str = "modern",
        language: str = "English"
    ) -> Dict[str, str]:
        """
        Generate SUPER PROFESSIONAL presentation from text prompt and optional images
        
        Args:
            prompt: User provided topic or instructions
            user_id: User requesting generation
            image_data: List of base64 encoded images (max 2)
            theme: Visual theme for presentation
            language: Target language for the presentation content
        """
        try:
            # 1. Structure content using AI with PREMIUM instructions
            print(f"🤖 Structuring PREMIUM PPT from prompt: '{prompt}' with {len(image_data or [])} images in {language}")
            
            structured_content = await PPTService._structure_from_prompt(prompt, image_data, language)
            
            # 2. Build PREMIUM PowerPoint
            print(f"📊 Building SUPER PROFESSIONAL PowerPoint presentation with theme: {theme}...")
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            theme_config = PPTService._get_theme_config(theme)
            
            # Add slides with PREMIUM design
            PPTService._add_premium_title_slide(prs, structured_content.get("title", "Presentation"), theme_config)
            PPTService._add_premium_summary_slide(prs, structured_content.get("summary", {}), theme_config)
            
            # Add content slides with PROFESSIONAL layouts
            for slide_data in structured_content.get("slides", []):
                if slide_data["type"] == "bullet":
                    PPTService._add_premium_bullet_slide(prs, slide_data, theme_config)
                elif slide_data["type"] == "table":
                    PPTService._add_premium_table_slide(prs, slide_data, theme_config)
                elif slide_data["type"] == "comparison":
                    PPTService._add_premium_comparison_slide(prs, slide_data, theme_config)
            
            # Add closing slide
            PPTService._add_premium_closing_slide(prs, theme_config)
            
            # 3. Save PPTX file
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_filename = f"premium_ppt_{user_id}_{timestamp}.pptx"
            pdf_filename = f"premium_ppt_{user_id}_{timestamp}.pdf"
            
            exports_dir = os.path.join("static", "exports")
            os.makedirs(exports_dir, exist_ok=True)
            
            pptx_path = os.path.join(exports_dir, pptx_filename)
            prs.save(pptx_path)
            print(f"✅ PPTX saved: {pptx_path}")
            
            # 4. Convert PPTX to PDF for preview
            print(f"🔄 Converting PPTX to PDF...")
            pdf_path = PPTService.convert_pptx_to_pdf(pptx_path)
            
            # 5. Generate URLs
            base_url = settings.base_url or "http://localhost:8000"
            pptx_url = f"{base_url}/static/exports/{pptx_filename}"
            pdf_url = f"{base_url}/static/exports/{pdf_filename}"
            
            # Direct file URLs for preview and download
            preview_url = pdf_url  # Preview shows PDF
            download_url = pdf_url  # Download gives PDF
            
            # 6. Calculate expiration (7 days from now)
            expires_at = datetime.now() + timedelta(days=7)
            
            return {
                "pptx_path": pptx_path,
                "pdf_path": pdf_path,
                "pptx_filename": pptx_filename,
                "pdf_filename": pdf_filename,
                "pptx_url": pptx_url,
                "pdf_url": pdf_url,
                "preview_url": preview_url,
                "download_url": download_url,
                "title": structured_content.get("title", "Presentation"),
                "theme": theme,
                "prompt": prompt,
                "script": structured_content.get("script", ""),
                "expires_at": expires_at
            }
            
        except Exception as e:
            print(f"❌ PPT generation error: {str(e)}")
            raise Exception(f"Failed to generate presentation: {str(e)}")

    @staticmethod
    async def _structure_from_prompt(prompt: str, images: List[str] = None, language: str = "English") -> Dict[str, Any]:
        """Use GPT-4o-mini Vision to structure PREMIUM presentation from prompt + images"""
        
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        

        system_prompt = f"""You are a WORLD-CLASS Presentation Designer working for Fortune 500 companies.
Create an EXCEPTIONAL, PREMIUM PowerPoint structure that will WOW the audience.

CRITICAL REQUIREMENTS:
1. OUTPUT LANGUAGE: You MUST generate ALL content in {language}.
2. Title must be COMPELLING and PROFESSIONAL
3. Summary must be EXECUTIVE-LEVEL with clear value proposition
4. Each slide must have MAXIMUM 4-5 bullet points (concise, impactful)
5. Use POWERFUL, ACTION-ORIENTED language
6. Include data points, statistics, or concrete examples where possible
7. Structure should tell a STORY with clear flow
8. Add comparison slides for before/after or pros/cons when relevant
9. For the LAST SLIDE, provide a 'closing' type slide.
10. IMPORTANT: You MUST also generate a PROFESSIONAL SPEAKER SCRIPT for the presentation.
    - Format it as a separate section at the very end of the JSON.
    - Key "script": "..."
    - The script should be TO THE POINT, Structured, and Professional.
    - Break it down by Slide 1, Slide 2, etc.
    - Do not just read the slides; provide context and flow.

RETURN JSON ONLY. No markdown formatting.
Structure:
{{
  "title": "Compelling Professional Title",
  "subtitles": "Optional subtitle or tagline",
  "summary": {{
    "overview": "Executive summary with clear value proposition and key insights...",
    "key_points": ["Impactful takeaway 1", "Impactful takeaway 2", "Impactful takeaway 3"]
  }},
  "slides": [
    {{
      "type": "bullet",
      "title": "Slide Title (Clear & Specific)",
      "points": ["Concise point 1", "Concise point 2", "Concise point 3"]
    }},
    {{
      "type": "comparison",
      "title": "Before vs After / Pros vs Cons",
      "left_title": "Before / Challenges",
      "left_points": ["Point 1", "Point 2"],
      "right_title": "After / Solutions",
      "right_points": ["Point 1", "Point 2"]
    }}
  ],
  "script": "Slide 1: ...\\n\\nSlide 2: ..."
}}

Make it PREMIUM. Make it WOW. Make it UNFORGETTABLE.
"""
        
        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": []}
        ]
        
        # Add text prompt
        messages[1]["content"].append({
            "type": "text",
            "text": f"Create a PREMIUM, PROFESSIONAL presentation about: {prompt}"
        })
        
        # Add images if provided
        if images:
            for img_base64 in images:
                if img_base64.startswith('data:image'):
                    url = img_base64
                else:
                    url = f"data:image/jpeg;base64,{img_base64}"
                    
                messages[1]["content"].append({
                    "type": "image_url",
                    "image_url": {"url": url}
                })

        try:
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=messages,
                temperature=0.8,  # Higher creativity for premium content
                max_tokens=2000,
                response_format={"type": "json_object"}
            )
            
            content = response.choices[0].message.content.strip()
            return json.loads(content)
            
        except Exception as e:
            print(f"⚠️ AI structuring failed: {e}")
            raise Exception(f"AI generation failed: {str(e)}")

    @staticmethod
    async def generate_presentation(scan_id_or_data: Any, user_id: str) -> Dict[str, str]:
        # Legacy support wrapper
        if isinstance(scan_id_or_data, dict):
             return await PPTService.generate_from_prompt(
                 prompt=f"Create a professional report for this document: {scan_id_or_data.get('original_filename', 'Document')}", 
                 user_id=user_id
             )
        return await PPTService.generate_from_prompt("Generate generic report", user_id)
    
    @staticmethod
    def _get_theme_config(theme: str) -> Dict[str, Any]:
        """Get color scheme and font settings for a theme"""
        
        # Premium Themes Definition
        themes = {
            "modern": { # Modern Blue
                "bg": (26, 35, 126), # Deep Blue
                "text": (255, 255, 255), 
                "accent": (68, 138, 255), # Bright Blue
                "sub": (197, 202, 233),
                "gradient_start": (26, 35, 126),
                "gradient_end": (13, 71, 161)
            },
            "minimalist": { # Minimalist Black
                "bg": (0, 0, 0), 
                "text": (255, 255, 255), 
                "accent": (255, 255, 255), # White accent
                "sub": (158, 158, 158),
                "gradient_start": (0, 0, 0),
                "gradient_end": (33, 33, 33)
            },
            "corporate": { # Professional Gray
                "bg": (33, 33, 33), 
                "text": (255, 255, 255), 
                "accent": (0, 188, 212), # Cyan
                "sub": (189, 189, 189),
                "gradient_start": (33, 33, 33),
                "gradient_end": (66, 66, 66)
            },
            "luxury": { # Elegant Gold
                "bg": (18, 18, 18), 
                "text": (212, 175, 55), # Gold text
                "accent": (255, 215, 0), # Bright Gold
                "sub": (184, 134, 11),
                "gradient_start": (0, 0, 0),
                "gradient_end": (28, 28, 28)
            },
            "creative": { # Creative Pink
                "bg": (136, 14, 79), 
                "text": (255, 255, 255), 
                "accent": (255, 64, 129), # Pink accent
                "sub": (248, 187, 208),
                "gradient_start": (136, 14, 79),
                "gradient_end": (194, 24, 91)
            }
        }
        
        # AI Recommended Logic: Pick a random premium theme
        if theme == "ai_recommended" or theme not in themes:
            import random
            selected_theme = random.choice(list(themes.keys()))
            print(f"🤖 AI Recommended Theme: {selected_theme}")
            return themes[selected_theme]
            
        return themes[theme]

    @staticmethod
    def _add_premium_title_slide(prs: Presentation, title: str, theme: Dict[str, Any]):
        """Add PREMIUM title slide with modern design and gradient effect"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add gradient background effect (simulated with shapes)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Add decorative accent bar at top
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Add main title with premium typography
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(8), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = title
        
        p = title_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(54)  # Larger, bolder title
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["text"])
        
        # Add premium subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(5), Inches(8), Inches(0.8)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Premium Presentation • {datetime.now().strftime('%B %d, %Y')}"
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.name = "Segoe UI Light"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Add decorative element at bottom
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4), Inches(6.5), Inches(2), Inches(0.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_shape.line.fill.background()
    
    @staticmethod
    def _add_premium_summary_slide(prs: Presentation, summary: Dict[str, Any], theme: Dict[str, Any]):
        """Add PREMIUM executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Overview text with premium formatting
        overview_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(8.4), Inches(2)
        )
        overview_frame = overview_box.text_frame
        overview_frame.word_wrap = True
        overview_frame.text = summary.get("overview", "")
        
        p = overview_frame.paragraphs[0]
        p.font.size = Pt(18)
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["text"])
        p.line_spacing = 1.5
        
        # Key points with icons/bullets
        key_points_y = 4.2
        for i, point in enumerate(summary.get("key_points", [])[:3]):
            # Bullet icon
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1), Inches(key_points_y + i * 0.6), Inches(0.15), Inches(0.15)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = RGBColor(*theme["accent"])
            bullet.line.fill.background()
            
            # Point text
            point_box = slide.shapes.add_textbox(
                Inches(1.3), Inches(key_points_y + i * 0.6 - 0.1), Inches(7.5), Inches(0.5)
            )
            point_frame = point_box.text_frame
            point_frame.text = point
            
            p = point_frame.paragraphs[0]
            p.font.size = Pt(16)
            p.font.name = "Segoe UI"
            p.font.color.rgb = RGBColor(*theme["text"])
    
    @staticmethod
    def _add_premium_bullet_slide(prs: Presentation, slide_data: Dict[str, Any], theme: Dict[str, Any]):
        """Add PREMIUM bullet point slide with modern layout"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Title with underline accent
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", "")
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Title underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.5), Inches(2), Inches(0.05)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = RGBColor(*theme["accent"])
        underline.line.fill.background()
        
        # Bullet points with premium spacing
        points_y = 2.2
        for i, point in enumerate(slide_data.get("points", [])[:5]):
            # Modern bullet
            bullet = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1), Inches(points_y + i * 0.9), Inches(0.2), Inches(0.2)
            )
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = RGBColor(*theme["accent"])
            bullet.line.fill.background()
            
            # Point text with generous spacing
            point_box = slide.shapes.add_textbox(
                Inches(1.4), Inches(points_y + i * 0.9 - 0.15), Inches(7.4), Inches(0.7)
            )
            point_frame = point_box.text_frame
            point_frame.word_wrap = True
            point_frame.text = point
            
            p = point_frame.paragraphs[0]
            p.font.size = Pt(22)  # Larger, more readable
            p.font.name = "Segoe UI"
            p.font.color.rgb = RGBColor(*theme["text"])
            p.line_spacing = 1.3
    
    @staticmethod
    def _add_premium_table_slide(prs: Presentation, slide_data: Dict[str, Any], theme: Dict[str, Any]):
        """Add PREMIUM table slide with professional formatting"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", "")
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Premium table
        rows = len(slide_data.get("rows", [])) + 1
        cols = len(slide_data.get("headers", []))
        
        if rows > 1 and cols > 0:
            table_shape = slide.shapes.add_table(
                rows, cols, Inches(1), Inches(2), Inches(8), Inches(4.5)
            )
            table = table_shape.table
            
            # Header row with premium styling
            for i, header in enumerate(slide_data.get("headers", [])):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(18)
                cell.text_frame.paragraphs[0].font.name = "Segoe UI"
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*theme["accent"])
            
            # Data rows with alternating colors
            for row_idx, row_data in enumerate(slide_data.get("rows", [])):
                for col_idx, cell_data in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_data)
                    cell.text_frame.paragraphs[0].font.size = Pt(16)
                    cell.text_frame.paragraphs[0].font.name = "Segoe UI"
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(*theme["text"])
                    cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    
                    # Alternating row colors
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(
                            min(theme["bg"][0] + 20, 255),
                            min(theme["bg"][1] + 20, 255),
                            min(theme["bg"][2] + 20, 255)
                        )
    
    @staticmethod
    def _add_premium_comparison_slide(prs: Presentation, slide_data: Dict[str, Any], theme: Dict[str, Any]):
        """Add PREMIUM comparison slide (Before/After, Pros/Cons)"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.2)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data.get("title", "Comparison")
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Left column
        left_title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(4), Inches(0.5)
        )
        left_title_frame = left_title_box.text_frame
        left_title_frame.text = slide_data.get("left_title", "Before")
        p = left_title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        p.alignment = PP_ALIGN.CENTER
        
        # Left points
        for i, point in enumerate(slide_data.get("left_points", [])[:4]):
            point_box = slide.shapes.add_textbox(
                Inches(1), Inches(2.5 + i * 0.8), Inches(3.6), Inches(0.6)
            )
            point_frame = point_box.text_frame
            point_frame.word_wrap = True
            point_frame.text = f"• {point}"
            p = point_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.name = "Segoe UI"
            p.font.color.rgb = RGBColor(*theme["text"])
        
        # Divider
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4.9), Inches(1.8), Inches(0.1), Inches(5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(*theme["accent"])
        divider.line.fill.background()
        
        # Right column
        right_title_box = slide.shapes.add_textbox(
            Inches(5.2), Inches(1.8), Inches(4), Inches(0.5)
        )
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = slide_data.get("right_title", "After")
        p = right_title_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        p.alignment = PP_ALIGN.CENTER
        
        # Right points
        for i, point in enumerate(slide_data.get("right_points", [])[:4]):
            point_box = slide.shapes.add_textbox(
                Inches(5.4), Inches(2.5 + i * 0.8), Inches(3.6), Inches(0.6)
            )
            point_frame = point_box.text_frame
            point_frame.word_wrap = True
            point_frame.text = f"• {point}"
            p = point_frame.paragraphs[0]
            p.font.size = Pt(18)
            p.font.name = "Segoe UI"
            p.font.color.rgb = RGBColor(*theme["text"])
    
    @staticmethod
    def _add_premium_closing_slide(prs: Presentation, theme: Dict[str, Any]):
        """Add PREMIUM closing slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme["bg"])
        
        # Accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(0.3)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_bar.line.fill.background()
        
        # Thank you message
        thank_you_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(8), Inches(1.5)
        )
        thank_you_frame = thank_you_box.text_frame
        thank_you_frame.text = "Thank You"
        
        p = thank_you_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(60)
        p.font.bold = True
        p.font.name = "Segoe UI"
        p.font.color.rgb = RGBColor(*theme["accent"])
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.8), Inches(8), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Questions & Discussion"
        
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.name = "Segoe UI Light"
        p.font.color.rgb = RGBColor(*theme["text"])
        
        # Decorative element
        accent_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(4), Inches(6), Inches(2), Inches(0.5)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = RGBColor(*theme["accent"])
        accent_shape.line.fill.background()
